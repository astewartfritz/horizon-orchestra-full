"""Horizon Orchestra — Document Generation Skill.

Create PDF, DOCX, PPTX, and XLSX documents from agent output.
Mirrors Perplexity's file-sharing capability: agents can hand users
polished documents, not just text.

Dependency matrix (graceful fallback if library not installed):
  PDF:   reportlab  →  weasyprint  →  wkhtmltopdf subprocess  →  .md
  DOCX:  python-docx  →  .md
  PPTX:  python-pptx  →  .md
  XLSX:  openpyxl  →  csv
"""

from __future__ import annotations

import asyncio
import csv
import logging
import os
import re
import tempfile
import time
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from .base import Skill

__all__ = ["DocumentGenerator"]

log = logging.getLogger("orchestra.skills.documents")

_WORKSPACE = Path(os.environ.get("ORCHESTRA_WORKSPACE", "/tmp/orchestra_docs"))


# ---------------------------------------------------------------------------
# Document generator
# ---------------------------------------------------------------------------

class DocumentGenerator(Skill):
    """Generate PDF, DOCX, PPTX, and XLSX documents from markdown content."""

    name: str = "documents"
    description: str = (
        "Create polished documents (PDF, DOCX, PPTX, XLSX) from markdown or "
        "structured data. Returns the file path of the saved document."
    )

    def __init__(self, workspace: str | Path | None = None) -> None:
        self.workspace = Path(workspace) if workspace else _WORKSPACE
        self.workspace.mkdir(parents=True, exist_ok=True)

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    async def create_pdf(
        self,
        content: str,
        title: str = "",
        output: str = "",
    ) -> str:
        """Convert markdown *content* to a styled PDF.

        Tries (in order):
        1. reportlab
        2. weasyprint
        3. wkhtmltopdf subprocess
        4. Save as .md fallback
        """
        output_path = self._ensure_output_path(output, title, "pdf")
        log.info("create_pdf() -> %s", output_path)

        # --- Try reportlab -----------------------------------------------
        try:
            return await asyncio.get_event_loop().run_in_executor(
                None, self._pdf_via_reportlab, content, title, output_path
            )
        except ImportError:
            log.debug("reportlab not available, trying weasyprint")
        except Exception as exc:
            log.warning("reportlab failed: %s", exc)

        # --- Try weasyprint ----------------------------------------------
        try:
            return await asyncio.get_event_loop().run_in_executor(
                None, self._pdf_via_weasyprint, content, title, output_path
            )
        except ImportError:
            log.debug("weasyprint not available, trying wkhtmltopdf")
        except Exception as exc:
            log.warning("weasyprint failed: %s", exc)

        # --- Try wkhtmltopdf / Chrome headless ---------------------------
        try:
            return await self._pdf_via_subprocess(content, title, output_path)
        except Exception as exc:
            log.warning("subprocess PDF failed: %s", exc)

        # --- Final fallback: save as markdown ----------------------------
        md_path = str(output_path).replace(".pdf", ".md")
        Path(md_path).write_text(f"# {title}\n\n{content}", encoding="utf-8")
        log.warning("PDF fallback: saved as markdown at %s", md_path)
        return md_path

    async def create_docx(
        self,
        content: str,
        title: str = "",
        output: str = "",
    ) -> str:
        """Convert markdown *content* to a DOCX file.

        Uses python-docx with proper heading, paragraph, code, and table styles.
        Falls back to saving as .md if python-docx is not available.
        """
        output_path = self._ensure_output_path(output, title, "docx")
        log.info("create_docx() -> %s", output_path)

        try:
            return await asyncio.get_event_loop().run_in_executor(
                None, self._docx_via_python_docx, content, title, output_path
            )
        except ImportError:
            log.warning("python-docx not available; saving as markdown")
        except Exception as exc:
            log.warning("DOCX creation failed: %s", exc)

        md_path = str(output_path).replace(".docx", ".md")
        Path(md_path).write_text(f"# {title}\n\n{content}", encoding="utf-8")
        return md_path

    async def create_pptx(
        self,
        content: str,
        title: str = "",
        slides_per_section: bool = True,
        output: str = "",
    ) -> str:
        """Convert markdown *content* to a PPTX presentation.

        Each ``##`` heading becomes a new slide.  Bullet lists become slide bullets.
        A title slide is prepended using *title*.

        Falls back to .md if python-pptx is not available.
        """
        output_path = self._ensure_output_path(output, title, "pptx")
        log.info("create_pptx() -> %s", output_path)

        try:
            return await asyncio.get_event_loop().run_in_executor(
                None, self._pptx_via_python_pptx, content, title, output_path
            )
        except ImportError:
            log.warning("python-pptx not available; saving as markdown")
        except Exception as exc:
            log.warning("PPTX creation failed: %s", exc)

        md_path = str(output_path).replace(".pptx", ".md")
        Path(md_path).write_text(f"# {title}\n\n{content}", encoding="utf-8")
        return md_path

    async def create_spreadsheet(
        self,
        data: list[dict[str, Any]],
        title: str = "",
        output: str = "",
    ) -> str:
        """Convert a list of dicts to an XLSX spreadsheet.

        Uses openpyxl with auto column widths and header styling.
        Falls back to CSV if openpyxl is not available.
        """
        output_path = self._ensure_output_path(output, title, "xlsx")
        log.info("create_spreadsheet() rows=%d -> %s", len(data), output_path)

        if not data:
            log.warning("No data provided for spreadsheet")
            Path(output_path).write_text("No data", encoding="utf-8")
            return str(output_path)

        try:
            return await asyncio.get_event_loop().run_in_executor(
                None, self._xlsx_via_openpyxl, data, title, output_path
            )
        except ImportError:
            log.warning("openpyxl not available; saving as CSV")
        except Exception as exc:
            log.warning("XLSX creation failed: %s", exc)

        # Fallback: CSV
        csv_path = str(output_path).replace(".xlsx", ".csv")
        headers = list(data[0].keys())
        with open(csv_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.DictWriter(f, fieldnames=headers)
            writer.writeheader()
            writer.writerows(data)
        return csv_path

    # ------------------------------------------------------------------
    # PDF backends
    # ------------------------------------------------------------------

    def _pdf_via_reportlab(self, content: str, title: str, output_path: str) -> str:
        """Generate PDF using reportlab."""
        from reportlab.lib.pagesizes import letter  # type: ignore[import]
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle  # type: ignore[import]
        from reportlab.lib.units import inch  # type: ignore[import]
        from reportlab.lib import colors  # type: ignore[import]
        from reportlab.platypus import (  # type: ignore[import]
            SimpleDocTemplate, Paragraph, Spacer, Preformatted
        )

        doc = SimpleDocTemplate(
            output_path,
            pagesize=letter,
            leftMargin=1 * inch,
            rightMargin=1 * inch,
            topMargin=1 * inch,
            bottomMargin=1 * inch,
        )
        styles = getSampleStyleSheet()
        story = []

        # Title
        if title:
            title_style = ParagraphStyle(
                "DocTitle",
                parent=styles["Heading1"],
                fontSize=24,
                spaceAfter=20,
                textColor=colors.HexColor("#1a1a2e"),
            )
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 0.2 * inch))

        sections = self._markdown_to_sections(content)
        for sec in sections:
            level = sec.get("level", 0)
            heading = sec.get("heading", "")
            body = sec.get("body", "")

            if heading:
                style_name = "Heading2" if level == 2 else ("Heading3" if level == 3 else "Heading1")
                story.append(Paragraph(heading, styles[style_name]))
                story.append(Spacer(1, 0.1 * inch))

            if body:
                # Detect code blocks
                code_pattern = re.compile(r"```[\w]*\n(.*?)```", re.DOTALL)
                last_end = 0
                for m in code_pattern.finditer(body):
                    pre_text = body[last_end:m.start()].strip()
                    if pre_text:
                        for para_text in pre_text.split("\n\n"):
                            para_text = para_text.strip()
                            if para_text:
                                story.append(Paragraph(para_text, styles["BodyText"]))
                                story.append(Spacer(1, 0.05 * inch))
                    code_text = m.group(1)
                    code_style = ParagraphStyle(
                        "Code",
                        parent=styles["Code"],
                        fontSize=8,
                        fontName="Courier",
                        backColor=colors.HexColor("#f5f5f5"),
                        leftIndent=20,
                    )
                    story.append(Preformatted(code_text, code_style))
                    story.append(Spacer(1, 0.1 * inch))
                    last_end = m.end()

                remaining = body[last_end:].strip()
                if remaining:
                    for para_text in remaining.split("\n\n"):
                        para_text = para_text.strip()
                        if para_text:
                            story.append(Paragraph(para_text, styles["BodyText"]))
                            story.append(Spacer(1, 0.05 * inch))

        doc.build(story)
        log.debug("reportlab PDF saved: %s", output_path)
        return output_path

    def _pdf_via_weasyprint(self, content: str, title: str, output_path: str) -> str:
        """Generate PDF using weasyprint (HTML→PDF)."""
        from weasyprint import HTML, CSS  # type: ignore[import]

        html_content = _markdown_to_html(content, title)
        css = CSS(string=_PDF_CSS)
        HTML(string=html_content).write_pdf(output_path, stylesheets=[css])
        log.debug("weasyprint PDF saved: %s", output_path)
        return output_path

    async def _pdf_via_subprocess(
        self, content: str, title: str, output_path: str
    ) -> str:
        """Generate PDF via wkhtmltopdf or Chrome headless."""
        html_content = _markdown_to_html(content, title)

        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".html", delete=False, encoding="utf-8"
        ) as f:
            f.write(html_content)
            html_path = f.name

        try:
            # Try wkhtmltopdf first
            proc = await asyncio.create_subprocess_exec(
                "wkhtmltopdf", "--quiet", html_path, output_path,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            _, stderr = await asyncio.wait_for(proc.communicate(), timeout=30)
            if proc.returncode == 0:
                log.debug("wkhtmltopdf PDF saved: %s", output_path)
                return output_path
            log.debug("wkhtmltopdf failed (rc=%d): %s", proc.returncode, stderr.decode())
        except (FileNotFoundError, asyncio.TimeoutError):
            pass

        # Try Chrome headless
        chrome_cmds = ["google-chrome", "chromium", "chromium-browser", "chrome"]
        try:
            for chrome in chrome_cmds:
                try:
                    proc = await asyncio.create_subprocess_exec(
                        chrome,
                        "--headless", "--no-sandbox", "--disable-gpu",
                        f"--print-to-pdf={output_path}",
                        html_path,
                        stdout=asyncio.subprocess.PIPE,
                        stderr=asyncio.subprocess.PIPE,
                    )
                    await asyncio.wait_for(proc.communicate(), timeout=30)
                    if proc.returncode == 0:
                        log.debug("Chrome headless PDF saved: %s", output_path)
                        return output_path
                except (FileNotFoundError, asyncio.TimeoutError):
                    continue

            raise RuntimeError("No PDF subprocess tool available")
        finally:
            Path(html_path).unlink(missing_ok=True)

    # ------------------------------------------------------------------
    # DOCX backend
    # ------------------------------------------------------------------

    def _docx_via_python_docx(self, content: str, title: str, output_path: str) -> str:
        """Generate DOCX using python-docx."""
        from docx import Document  # type: ignore[import]
        from docx.shared import Pt, RGBColor  # type: ignore[import]
        from docx.enum.text import WD_ALIGN_PARAGRAPH  # type: ignore[import]

        doc = Document()

        # Set up normal style
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)

        # Title
        if title:
            title_para = doc.add_heading(title, level=0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        sections = self._markdown_to_sections(content)
        for sec in sections:
            level = sec.get("level", 0)
            heading = sec.get("heading", "")
            body = sec.get("body", "")

            if heading and level > 0:
                doc.add_heading(heading, level=min(level, 9))

            if body:
                code_pattern = re.compile(r"```[\w]*\n(.*?)```", re.DOTALL)
                last_end = 0
                for m in code_pattern.finditer(body):
                    pre_text = body[last_end:m.start()].strip()
                    if pre_text:
                        _add_paragraphs(doc, pre_text)
                    # Code block — use monospace
                    code_para = doc.add_paragraph(m.group(1))
                    code_para.style = "No Spacing"
                    code_para.runs[0].font.name = "Courier New"
                    code_para.runs[0].font.size = Pt(9)
                    last_end = m.end()

                remaining = body[last_end:].strip()
                if remaining:
                    _add_paragraphs(doc, remaining)

        doc.save(output_path)
        log.debug("python-docx DOCX saved: %s", output_path)
        return output_path

    # ------------------------------------------------------------------
    # PPTX backend
    # ------------------------------------------------------------------

    def _pptx_via_python_pptx(self, content: str, title: str, output_path: str) -> str:
        """Generate PPTX using python-pptx."""
        from pptx import Presentation  # type: ignore[import]
        from pptx.util import Inches, Pt  # type: ignore[import]
        from pptx.dml.color import RGBColor  # type: ignore[import]
        from pptx.enum.text import PP_ALIGN  # type: ignore[import]

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]   # blank
        title_layout = prs.slide_layouts[0]   # title slide

        # --- Title slide ------------------------------------------------
        if title:
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = title
            if slide.placeholders[1]:
                slide.placeholders[1].text = ""

        # --- Content slides from sections -------------------------------
        sections = self._markdown_to_sections(content)

        for sec in sections:
            level = sec.get("level", 0)
            heading = sec.get("heading", "")
            body = sec.get("body", "")

            if level > 2:
                # Sub-sections: append to previous slide's body if possible
                continue

            slide = prs.slides.add_slide(blank_layout)

            # Heading box
            if heading:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = heading
                p.font.bold = True
                p.font.size = Pt(28)
                p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

            # Body box
            if body:
                txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12), Inches(5.5))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                # Strip code fences for slides
                clean_body = re.sub(r"```[\w]*\n.*?```", "[code block]", body, flags=re.DOTALL)
                lines = clean_body.strip().splitlines()
                first = True
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue
                    if first:
                        para = tf2.paragraphs[0]
                        first = False
                    else:
                        para = tf2.add_paragraph()
                    # Bullet items
                    if line.startswith("- ") or line.startswith("* "):
                        para.text = "•  " + line[2:]
                        para.font.size = Pt(18)
                    elif re.match(r"^\d+\.\s", line):
                        para.text = line
                        para.font.size = Pt(18)
                    else:
                        para.text = line
                        para.font.size = Pt(16)

        prs.save(output_path)
        log.debug("python-pptx PPTX saved: %s", output_path)
        return output_path

    # ------------------------------------------------------------------
    # XLSX backend
    # ------------------------------------------------------------------

    def _xlsx_via_openpyxl(
        self, data: list[dict[str, Any]], title: str, output_path: str
    ) -> str:
        """Generate XLSX using openpyxl with auto-widths and header styling."""
        from openpyxl import Workbook  # type: ignore[import]
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side  # type: ignore[import]
        from openpyxl.utils import get_column_letter  # type: ignore[import]

        wb = Workbook()
        ws = wb.active
        ws.title = title[:31] if title else "Sheet1"  # max 31 chars

        headers = list(data[0].keys())

        # Header row styling
        header_font = Font(bold=True, color="FFFFFF", name="Calibri", size=11)
        header_fill = PatternFill(start_color="1A1A2E", end_color="1A1A2E", fill_type="solid")
        header_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
        thin_border = Border(
            left=Side(style="thin"),
            right=Side(style="thin"),
            top=Side(style="thin"),
            bottom=Side(style="thin"),
        )

        for col_idx, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col_idx, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_align
            cell.border = thin_border

        # Data rows
        alt_fill = PatternFill(start_color="F0F4FF", end_color="F0F4FF", fill_type="solid")
        for row_idx, row in enumerate(data, 2):
            fill = alt_fill if row_idx % 2 == 0 else None
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=row.get(header, ""))
                cell.border = thin_border
                if fill:
                    cell.fill = fill

        # Auto-width columns
        for col_idx, header in enumerate(headers, 1):
            max_len = len(str(header))
            for row in data:
                val = str(row.get(header, ""))
                max_len = max(max_len, len(val))
            ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len + 4, 60)

        ws.freeze_panes = "A2"

        wb.save(output_path)
        log.debug("openpyxl XLSX saved: %s", output_path)
        return output_path

    # ------------------------------------------------------------------
    # Markdown parser
    # ------------------------------------------------------------------

    def _markdown_to_sections(self, content: str) -> list[dict[str, Any]]:
        """Parse markdown into sections.

        Returns a list of dicts: ``{heading, body, level}``.
        Level 0 means preamble (before any heading).
        """
        sections: list[dict[str, Any]] = []
        current: dict[str, Any] = {"heading": "", "body": "", "level": 0}
        lines = content.splitlines()
        in_code_block = False

        for line in lines:
            # Track code fences so we don't misparse headings inside them
            stripped = line.strip()
            if stripped.startswith("```"):
                in_code_block = not in_code_block

            if not in_code_block:
                m = re.match(r"^(#{1,6})\s+(.*)", line)
                if m:
                    # Save current section
                    if current["heading"] or current["body"].strip():
                        sections.append(current)
                    level = len(m.group(1))
                    current = {"heading": m.group(2).strip(), "body": "", "level": level}
                    continue

            current["body"] += line + "\n"

        if current["heading"] or current["body"].strip():
            sections.append(current)

        return sections

    # ------------------------------------------------------------------
    # Utility
    # ------------------------------------------------------------------

    def _ensure_output_path(self, output: str, title: str, ext: str) -> str:
        """Resolve or generate an output file path."""
        if output:
            p = Path(output)
            p.parent.mkdir(parents=True, exist_ok=True)
            return str(p)
        slug = re.sub(r"[^\w\-]", "_", title)[:40] if title else str(uuid.uuid4())[:8]
        ts = int(time.time())
        filename = f"{slug}_{ts}.{ext}"
        return str(self.workspace / filename)

    def _ensure_output_dir(self, path: str) -> str:
        """Ensure parent dir of *path* exists."""
        Path(path).parent.mkdir(parents=True, exist_ok=True)
        return path

    # ------------------------------------------------------------------
    # Skill ABC interface
    # ------------------------------------------------------------------

    def get_tool_definitions(self) -> list[dict[str, Any]]:
        return [
            {
                "type": "function",
                "function": {
                    "name": "doc_create_pdf",
                    "description": "Convert markdown content to a styled PDF file. Returns file path.",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "content": {"type": "string", "description": "Markdown content."},
                            "title": {"type": "string", "description": "Document title.", "default": ""},
                            "output": {"type": "string", "description": "Output file path (optional).", "default": ""},
                        },
                        "required": ["content"],
                    },
                },
            },
            {
                "type": "function",
                "function": {
                    "name": "doc_create_docx",
                    "description": "Convert markdown content to a DOCX Word document. Returns file path.",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "content": {"type": "string", "description": "Markdown content."},
                            "title": {"type": "string", "description": "Document title.", "default": ""},
                            "output": {"type": "string", "description": "Output file path (optional).", "default": ""},
                        },
                        "required": ["content"],
                    },
                },
            },
            {
                "type": "function",
                "function": {
                    "name": "doc_create_pptx",
                    "description": (
                        "Convert markdown content to a PPTX presentation. "
                        "Each ## heading becomes a slide. Returns file path."
                    ),
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "content": {"type": "string", "description": "Markdown content."},
                            "title": {"type": "string", "description": "Presentation title.", "default": ""},
                            "output": {"type": "string", "description": "Output file path (optional).", "default": ""},
                        },
                        "required": ["content"],
                    },
                },
            },
            {
                "type": "function",
                "function": {
                    "name": "doc_create_xlsx",
                    "description": "Convert a list of records to an XLSX spreadsheet. Returns file path.",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "data": {
                                "type": "array",
                                "items": {"type": "object"},
                                "description": "List of row dicts.",
                            },
                            "title": {"type": "string", "description": "Sheet title.", "default": ""},
                            "output": {"type": "string", "description": "Output file path (optional).", "default": ""},
                        },
                        "required": ["data"],
                    },
                },
            },
        ]

    async def execute(self, action: str, params: dict[str, Any]) -> dict[str, Any]:
        """Dispatch tool calls."""
        if action == "doc_create_pdf":
            path = await self.create_pdf(
                content=params["content"],
                title=params.get("title", ""),
                output=params.get("output", ""),
            )
            return {"path": path, "format": "pdf"}

        if action == "doc_create_docx":
            path = await self.create_docx(
                content=params["content"],
                title=params.get("title", ""),
                output=params.get("output", ""),
            )
            return {"path": path, "format": "docx"}

        if action == "doc_create_pptx":
            path = await self.create_pptx(
                content=params["content"],
                title=params.get("title", ""),
                output=params.get("output", ""),
            )
            return {"path": path, "format": "pptx"}

        if action == "doc_create_xlsx":
            path = await self.create_spreadsheet(
                data=params["data"],
                title=params.get("title", ""),
                output=params.get("output", ""),
            )
            return {"path": path, "format": "xlsx"}

        return {"error": f"Unknown action: {action!r}"}


# ---------------------------------------------------------------------------
# Module helpers
# ---------------------------------------------------------------------------

def _add_paragraphs(doc: Any, text: str) -> None:
    """Add paragraphs to a python-docx Document, handling bullet lists."""
    for block in text.split("\n\n"):
        block = block.strip()
        if not block:
            continue
        lines = block.splitlines()
        for line in lines:
            line = line.strip()
            if not line:
                continue
            if line.startswith("- ") or line.startswith("* "):
                p = doc.add_paragraph(line[2:], style="List Bullet")
            elif re.match(r"^\d+\.\s", line):
                p = doc.add_paragraph(re.sub(r"^\d+\.\s", "", line), style="List Number")
            else:
                p = doc.add_paragraph(line)


def _markdown_to_html(content: str, title: str) -> str:
    """Convert markdown to a minimal HTML document for PDF conversion."""
    # Try markdown library; fall back to basic conversion
    try:
        import markdown as md_lib  # type: ignore[import]
        body_html = md_lib.markdown(
            content, extensions=["tables", "fenced_code", "codehilite"]
        )
    except ImportError:
        # Basic conversion
        body_html = content
        body_html = re.sub(r"^# (.+)$", r"<h1>\1</h1>", body_html, flags=re.MULTILINE)
        body_html = re.sub(r"^## (.+)$", r"<h2>\1</h2>", body_html, flags=re.MULTILINE)
        body_html = re.sub(r"^### (.+)$", r"<h3>\1</h3>", body_html, flags=re.MULTILINE)
        body_html = re.sub(r"```[\w]*\n(.*?)```", r"<pre><code>\1</code></pre>", body_html, flags=re.DOTALL)
        body_html = re.sub(r"`([^`]+)`", r"<code>\1</code>", body_html)
        body_html = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", body_html)
        body_html = re.sub(r"\*(.+?)\*", r"<em>\1</em>", body_html)
        body_html = "<p>" + re.sub(r"\n\n+", "</p><p>", body_html) + "</p>"

    return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>{title}</title>
<style>{_PDF_CSS}</style>
</head>
<body>
{'<h1>' + title + '</h1>' if title else ''}
{body_html}
</body>
</html>"""


_PDF_CSS = """
body { font-family: 'Helvetica Neue', Arial, sans-serif; font-size: 12pt;
       line-height: 1.6; color: #1a1a2e; max-width: 800px; margin: 40px auto; padding: 0 20px; }
h1 { font-size: 24pt; color: #1a1a2e; border-bottom: 2px solid #4a90d9; padding-bottom: 8px; }
h2 { font-size: 18pt; color: #2d5a8e; margin-top: 24px; }
h3 { font-size: 14pt; color: #4a7fb5; }
pre { background: #f5f7ff; border-left: 4px solid #4a90d9; padding: 12px 16px;
      font-family: 'Courier New', monospace; font-size: 10pt; overflow-x: auto; }
code { font-family: 'Courier New', monospace; font-size: 10pt;
       background: #f0f4ff; padding: 2px 5px; border-radius: 3px; }
table { border-collapse: collapse; width: 100%; margin: 16px 0; }
th { background: #1a1a2e; color: white; padding: 8px 12px; text-align: left; }
td { padding: 8px 12px; border-bottom: 1px solid #dde; }
tr:nth-child(even) { background: #f5f7ff; }
a { color: #4a90d9; }
"""
