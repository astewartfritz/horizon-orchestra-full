"""Horizon Orchestra — PowerPoint Generation.

Create presentations with multiple slide layouts, themes, images,
charts, and custom styling.  Uses ``python-pptx`` with a try/except
guard.

Usage::

    from orchestra.documents.pptx import PPTXGenerator, SlideSpec

    gen = PPTXGenerator()
    slides = [
        SlideSpec(title="Welcome", content="Introduction", layout="title_slide"),
        SlideSpec(title="Data", bullets=["Point 1", "Point 2"], layout="content"),
    ]
    pptx_bytes = gen.create("My Presentation", slides)
"""

from __future__ import annotations

import io
import logging
import os
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Optional, Sequence

__all__ = [
    "PPTXGenerator",
    "SlideSpec",
]

log = logging.getLogger("orchestra.documents.pptx")

_WORKSPACE = Path(os.environ.get("ORCHESTRA_WORKSPACE", "/tmp/orchestra_docs"))

# Optional dependency: python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE
    _HAS_PPTX = True
except ImportError:
    Presentation = None  # type: ignore[assignment, misc]
    _HAS_PPTX = False


# ---------------------------------------------------------------------------
# Dataclasses
# ---------------------------------------------------------------------------

@dataclass
class ThemeConfig:
    """Color and font theme for presentations."""

    primary_color: str = "1E3A5F"      # Dark navy
    secondary_color: str = "3B82F6"    # Blue
    accent_color: str = "10B981"       # Green
    background_color: str = "FFFFFF"   # White
    text_color: str = "1A1A1A"         # Near-black
    light_text: str = "FFFFFF"         # White
    heading_font: str = "Calibri"
    body_font: str = "Calibri"
    heading_size: int = 32
    body_size: int = 18
    subtitle_size: int = 24


@dataclass
class SlideSpec:
    """Specification for a single slide.

    Parameters
    ----------
    title:
        Slide title text.
    content:
        Main body text.
    bullets:
        Bullet point list.
    image:
        Path to an image file to embed.
    chart:
        Chart data dict (``type``, ``categories``, ``series``).
    layout:
        Layout name (``title_slide``, ``section_header``, ``content``,
        ``content_with_image``, ``two_column``, ``chart_slide``, ``blank``).
    notes:
        Speaker notes text.
    """

    title: str = ""
    content: str = ""
    bullets: list[str] = field(default_factory=list)
    image: str = ""
    chart: dict[str, Any] = field(default_factory=dict)
    layout: str = "content"
    notes: str = ""
    subtitle: str = ""
    left_content: str = ""
    right_content: str = ""
    left_bullets: list[str] = field(default_factory=list)
    right_bullets: list[str] = field(default_factory=list)


# ---------------------------------------------------------------------------
# PPTXGenerator
# ---------------------------------------------------------------------------

class PPTXGenerator:
    """PowerPoint presentation generator using python-pptx.

    Parameters
    ----------
    workspace:
        Directory for saving output files.
    theme:
        Color and font theme configuration.
    """

    LAYOUTS = {
        "title_slide",
        "section_header",
        "content",
        "content_with_image",
        "two_column",
        "chart_slide",
        "blank",
    }

    def __init__(
        self,
        workspace: str | Path | None = None,
        theme: ThemeConfig | None = None,
    ) -> None:
        if not _HAS_PPTX:
            raise ImportError(
                "python-pptx is required for PowerPoint generation: "
                "pip install python-pptx"
            )
        self.workspace = Path(workspace) if workspace else _WORKSPACE / "pptx"
        self.workspace.mkdir(parents=True, exist_ok=True)
        self.theme = theme or ThemeConfig()

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _save_path(self) -> Path:
        return self.workspace / f"{uuid.uuid4().hex[:12]}.pptx"

    def _rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color string to RGBColor."""
        hex_color = hex_color.lstrip("#")
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )

    def _set_text_style(
        self,
        run: Any,
        *,
        size: int = 18,
        bold: bool = False,
        color: str = "",
        font_name: str = "",
    ) -> None:
        """Apply styling to a text run."""
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = self._rgb(color)
        if font_name:
            run.font.name = font_name

    def _add_title(self, slide: Any, title: str) -> None:
        """Add a styled title to a slide."""
        if not title:
            return
        if slide.shapes.title:
            tf = slide.shapes.title.text_frame
            tf.text = title
            for para in tf.paragraphs:
                for run in para.runs:
                    self._set_text_style(
                        run,
                        size=self.theme.heading_size,
                        bold=True,
                        color=self.theme.primary_color,
                        font_name=self.theme.heading_font,
                    )

    def _add_body_text(self, text_frame: Any, text: str) -> None:
        """Add body text to a text frame."""
        para = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        para.text = text
        for run in para.runs:
            self._set_text_style(
                run,
                size=self.theme.body_size,
                color=self.theme.text_color,
                font_name=self.theme.body_font,
            )

    def _add_bullets(self, text_frame: Any, bullets: list[str]) -> None:
        """Add bullet points to a text frame."""
        for i, bullet in enumerate(bullets):
            if i == 0 and not text_frame.paragraphs[0].text:
                para = text_frame.paragraphs[0]
            else:
                para = text_frame.add_paragraph()
            para.text = bullet
            para.level = 0
            para.space_after = Pt(6)
            for run in para.runs:
                self._set_text_style(
                    run,
                    size=self.theme.body_size,
                    color=self.theme.text_color,
                    font_name=self.theme.body_font,
                )

    def _get_layout_index(self, prs: Any, layout_name: str) -> int:
        """Map layout name to slide layout index."""
        layout_map = {
            "title_slide": 0,
            "section_header": 2,
            "content": 1,
            "content_with_image": 1,
            "two_column": 3,
            "chart_slide": 1,
            "blank": 6,
        }
        idx = layout_map.get(layout_name, 1)
        # Clamp to available layouts
        max_idx = len(prs.slide_layouts) - 1
        return min(idx, max_idx)

    # ------------------------------------------------------------------
    # Slide builders
    # ------------------------------------------------------------------

    def _build_title_slide(self, prs: Any, spec: SlideSpec) -> Any:
        """Build a title slide."""
        layout_idx = self._get_layout_index(prs, "title_slide")
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        self._add_title(slide, spec.title)

        # Subtitle
        if spec.subtitle or spec.content:
            subtitle_text = spec.subtitle or spec.content
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    shape.text = subtitle_text
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            self._set_text_style(
                                run,
                                size=self.theme.subtitle_size,
                                color=self.theme.secondary_color,
                                font_name=self.theme.body_font,
                            )
                    break
        return slide

    def _build_section_header(self, prs: Any, spec: SlideSpec) -> Any:
        """Build a section header slide."""
        layout_idx = self._get_layout_index(prs, "section_header")
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        self._add_title(slide, spec.title)

        if spec.content:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    shape.text = spec.content
                    break
        return slide

    def _build_content_slide(self, prs: Any, spec: SlideSpec) -> Any:
        """Build a content slide with text and/or bullets."""
        layout_idx = self._get_layout_index(prs, "content")
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        self._add_title(slide, spec.title)

        # Find the body placeholder
        body_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_shape = shape
                break

        if body_shape:
            tf = body_shape.text_frame
            tf.word_wrap = True

            if spec.content:
                self._add_body_text(tf, spec.content)

            if spec.bullets:
                if spec.content:
                    # Add spacing after content
                    tf.add_paragraph()
                self._add_bullets(tf, spec.bullets)

        return slide

    def _build_content_with_image(self, prs: Any, spec: SlideSpec) -> Any:
        """Build a slide with content on the left and an image on the right."""
        layout_idx = self._get_layout_index(prs, "content")
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        self._add_title(slide, spec.title)

        # Add text in left half
        from pptx.util import Inches
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(4.5)
        height = Inches(4.5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        if spec.content:
            self._add_body_text(tf, spec.content)
        if spec.bullets:
            self._add_bullets(tf, spec.bullets)

        # Add image on right
        if spec.image and Path(spec.image).exists():
            img_left = Inches(5.5)
            img_top = Inches(1.8)
            img_width = Inches(4.0)
            slide.shapes.add_picture(
                spec.image, img_left, img_top, width=img_width,
            )

        return slide

    def _build_two_column(self, prs: Any, spec: SlideSpec) -> Any:
        """Build a two-column slide."""
        layout_idx = self._get_layout_index(prs, "content")
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        self._add_title(slide, spec.title)

        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(4.3), Inches(4.5),
        )
        left_tf = left_box.text_frame
        left_tf.word_wrap = True

        if spec.left_content:
            self._add_body_text(left_tf, spec.left_content)
        if spec.left_bullets:
            self._add_bullets(left_tf, spec.left_bullets)

        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(5.3), Inches(1.8), Inches(4.3), Inches(4.5),
        )
        right_tf = right_box.text_frame
        right_tf.word_wrap = True

        if spec.right_content:
            self._add_body_text(right_tf, spec.right_content)
        if spec.right_bullets:
            self._add_bullets(right_tf, spec.right_bullets)

        return slide

    def _build_chart_slide(self, prs: Any, spec: SlideSpec) -> Any:
        """Build a slide with a chart."""
        layout_idx = self._get_layout_index(prs, "content")
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        self._add_title(slide, spec.title)

        if not spec.chart:
            return slide

        chart_type_str = spec.chart.get("type", "bar")
        categories = spec.chart.get("categories", [])
        series_data = spec.chart.get("series", [])

        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "area": XL_CHART_TYPE.AREA,
        }
        xl_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        from pptx.chart.data import CategoryChartData
        chart_data = CategoryChartData()
        chart_data.categories = categories

        for s in series_data:
            chart_data.add_series(s.get("name", "Series"), s.get("values", []))

        chart_frame = slide.shapes.add_chart(
            xl_type,
            Inches(1.0), Inches(1.8),
            Inches(8.0), Inches(5.0),
            chart_data,
        )

        # Style the chart
        chart = chart_frame.chart
        chart.has_legend = len(series_data) > 1
        if chart.has_legend:
            chart.legend.include_in_layout = False

        return slide

    def _build_blank_slide(self, prs: Any, spec: SlideSpec) -> Any:
        """Build a blank slide."""
        layout_idx = self._get_layout_index(prs, "blank")
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if spec.title:
            txBox = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9.0), Inches(1.0),
            )
            tf = txBox.text_frame
            para = tf.paragraphs[0]
            para.text = spec.title
            for run in para.runs:
                self._set_text_style(
                    run,
                    size=self.theme.heading_size,
                    bold=True,
                    color=self.theme.primary_color,
                    font_name=self.theme.heading_font,
                )

        return slide

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def create(
        self,
        title: str,
        slides: Sequence[SlideSpec],
        *,
        theme: ThemeConfig | None = None,
    ) -> bytes:
        """Create a complete PowerPoint presentation.

        Parameters
        ----------
        title:
            Presentation title (used for the first slide if no title_slide
            spec is provided).
        slides:
            List of slide specifications.
        theme:
            Optional theme override.

        Returns
        -------
        bytes
            PPTX file content.
        """
        if theme:
            self.theme = theme

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        builders = {
            "title_slide": self._build_title_slide,
            "section_header": self._build_section_header,
            "content": self._build_content_slide,
            "content_with_image": self._build_content_with_image,
            "two_column": self._build_two_column,
            "chart_slide": self._build_chart_slide,
            "blank": self._build_blank_slide,
        }

        for spec in slides:
            layout = spec.layout if spec.layout in builders else "content"
            builder = builders[layout]
            slide = builder(prs, spec)

            # Add speaker notes
            if spec.notes and slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                notes_tf.text = spec.notes

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def add_slide(
        self,
        pptx_data: bytes,
        spec: SlideSpec,
    ) -> bytes:
        """Add a single slide to an existing presentation.

        Parameters
        ----------
        pptx_data:
            Existing PPTX file content.
        spec:
            Slide specification.

        Returns
        -------
        bytes
            Updated PPTX content.
        """
        prs = Presentation(io.BytesIO(pptx_data))

        builders = {
            "title_slide": self._build_title_slide,
            "section_header": self._build_section_header,
            "content": self._build_content_slide,
            "content_with_image": self._build_content_with_image,
            "two_column": self._build_two_column,
            "chart_slide": self._build_chart_slide,
            "blank": self._build_blank_slide,
        }

        layout = spec.layout if spec.layout in builders else "content"
        builder = builders[layout]
        slide = builder(prs, spec)

        if spec.notes and slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = spec.notes

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def save(self, pptx_bytes: bytes, path: str | Path | None = None) -> Path:
        """Save PPTX bytes to a file.

        Parameters
        ----------
        pptx_bytes:
            PPTX content.
        path:
            Output path.

        Returns
        -------
        Path
            Saved file path.
        """
        save_to = Path(path) if path else self._save_path()
        save_to.parent.mkdir(parents=True, exist_ok=True)
        save_to.write_bytes(pptx_bytes)
        log.info("Saved PPTX → %s (%d bytes)", save_to, len(pptx_bytes))
        return save_to
